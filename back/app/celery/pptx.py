from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

def create_prs(
    prs,
    slides,
):
    # Центрируем заголовок
    # title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    alignment_map = {
        'center': PP_ALIGN.CENTER,
        'left': PP_ALIGN.LEFT,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }

    for slide_item in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for item in slide_item['elements']:
            main_text_shape = slide.shapes.add_textbox(
                item["x"]*prs.slide_width,
                item["y"]*prs.slide_height,
                item["x"]*prs.slide_width+item["w"]*prs.slide_width,
                item["y"]*prs.slide_height+item["h"]*prs.slide_height
            )
            main_text_tf = main_text_shape.text_frame
            main_text_tf.text = item['content']
            main_text_p = main_text_tf.paragraphs[0]
            main_text_p.alignment = alignment_map[item["alignment"]]
            main_text_run = main_text_p.runs[0]
            main_text_run.font.size = Pt(item['size'])
            if item['style'] == 'bold':
                main_text_run.font.bold = True
            if item['style'] == 'italic':
                main_text_run.font.italic = True
    
    # main_text_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    return prs


def get_pres(pres_json, id, design_number=1):
    prs = Presentation(f"pptx/{design_number}.pptx")
    create_prs(prs, pres_json['slides'])
    prs.save(f"{id}.pptx")

# pres = {
#   "slides": [
#     {
#       "id": "550e8400-e29b-41d4-a716-446655440000",
#       "slide_number": 1,
#       "elements": [
#         {
#           "text_type": "header", # regular, header, list
#           "alignment": "center", # left, right, center, justify
#           "style": "bold", # regular, bold, italic
#           "size": 48,
#           "content": "Заголовок презентации",
#           "coords": [0.2,0.2,0.4,0.4]
#         },
#         {
#           "text_type": "regular",
#           "alignment": "left",
#           "style": "regular",
#           "size": 24,
#           "content": "Обычный текст на слайде",
#           "coords": [0.5,0.5,0.7,0.7]
#         }
#       ]
#     }
#   ]
# }